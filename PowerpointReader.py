from pptx import Presentation
from pptx.slide import Slide

import SlideData


class PowerpointReader:
    def __init__(self, _path):
        self.powerpoint = Presentation(_path)

    def read_slides(self) -> list[SlideData.SlideData]:
        slides_data = []
        for i, slide in enumerate(self.powerpoint.slides):
            slides_data.append(self.extract_data(slide))
        return slides_data

    def extract_data(self, slide: Slide) -> SlideData.SlideData:
        titles, body = [], []
        for shape in slide.placeholders:
            if shape.placeholder_format.type in [1, 3]:
                titles.append(shape.text)
            elif shape.placeholder_format.type == 2:
                body.append(shape.text)
            elif shape.placeholder_format.type == 7 and shape.has_text_frame and shape.text:
                body.append(shape.text)
        return SlideData.SlideData(self.powerpoint.core_properties.title, titles, body)
